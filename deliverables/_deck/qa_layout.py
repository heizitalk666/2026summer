# -*- coding: utf-8 -*-
"""版面 QA：出界 / 底部留白 / 实体元素重叠 / 表格与文本框的换行溢出估算。

pptxgenjs 写给表格 graphicFrame 的 height 是名义值（本文件里恒为 1.0"），
真实高度要按行高求和；渲染器还会按内容再撑高行，所以这里同时做一次换行估算。
模板套版（顶部横幅、校徽、标题块、底条）与整版封面页按设计豁免。
"""
import math, sys
from pptx import Presentation
from pptx.util import Emu

PATH = sys.argv[1] if len(sys.argv) > 1 else "中期答辩.pptx"
p = Presentation(PATH)
SH, SW = Emu(p.slide_height).inches, Emu(p.slide_width).inches
FULL_BLEED = {1, len(p.slides)}
PAD = 0.10          # 单元格左右内边距合计
LS_CELL = 1.60      # 表格单元格：实测 LibreOffice 渲染 10.5pt 中文表头每行约占 0.24"，含上下边距
LS_BOX  = 1.28      # 文本框：margin 0，行距接近字号本身

def width_units(txt):
    """全角字按 1，半角按 0.55 折算成「字宽」数。"""
    u = 0.0
    for ch in txt:
        u += 1.0 if ord(ch) > 0x2E80 else 0.55
    return u

def lines_needed(txt, box_w_in, pt):
    if not txt.strip():
        return 0
    cap = max(1.0, (box_w_in - PAD) * 72.0 / pt)      # 每行可容纳的字宽数
    n = 0
    for seg in txt.split("\n"):
        n += max(1, math.ceil(width_units(seg) / cap))
    return n

def text_pt(cell_or_tf, default=10.0):
    for para in cell_or_tf.paragraphs:
        for r in para.runs:
            if r.font.size:
                return r.font.size.pt
    return default

def table_geom(sh):
    """返回 (left, top, right, bottom_estimated)。"""
    t = sh.table
    left, top = Emu(sh.left).inches, Emu(sh.top).inches
    colw = [Emu(c.width).inches for c in t.columns]
    h = 0.0
    for r in t.rows:
        declared = Emu(r.height).inches
        need = 0.0
        for ci, cell in enumerate(r.cells):
            pt = text_pt(cell.text_frame)
            n = lines_needed(cell.text, colw[ci], pt)
            need = max(need, n * pt * LS_CELL / 72.0 + 0.10)
        h += max(declared, need)
    # 行高估算仍会比渲染结果偏小若干个百分点，留 0.12" 安全余量
    return left, top, left + sum(colw), top + h + 0.12

def geom(sh):
    if getattr(sh, "has_table", False) and sh.has_table:
        return table_geom(sh)
    l, t = Emu(sh.left).inches, Emu(sh.top).inches
    return l, t, l + Emu(sh.width).inches, t + Emu(sh.height).inches

def solid(sh):
    return str(sh.shape_type) in ("TABLE (19)", "PICTURE (13)", "AUTO_SHAPE (1)")

bad = []
for i, s in enumerate(p.slides, 1):
    if i in FULL_BLEED:
        continue
    items, bot, over, spill = [], 0.0, [], []
    for sh in s.shapes:
        try:
            b = geom(sh)
        except (TypeError, AttributeError):
            continue
        if b[3] <= 0.92 or b[1] >= SH - 0.30:          # 页眉套版 / 底条
            continue
        # 文本框换行溢出估算
        if sh.has_text_frame and sh.text_frame.text.strip():
            pt = text_pt(sh.text_frame)
            need = lines_needed(sh.text_frame.text, b[2] - b[0], pt) * pt * LS_BOX / 72.0
            declared = b[3] - b[1]
            if need > declared + 0.06:
                spill.append((sh.text_frame.text[:14].replace("\n", "/"),
                              round(need, 2), round(declared, 2)))
        if b[1] > SH - 0.75:
            continue
        bot = max(bot, b[3])
        if b[3] > SH - 0.30 or b[2] > SW + 0.01 or b[0] < -0.01:
            over.append((str(sh.shape_type).split()[0], round(b[0], 2), round(b[2], 2), round(b[3], 2)))
        if solid(sh):
            items.append((b, str(sh.shape_type)))
    ov = []
    for a in range(len(items)):
        for c in range(a + 1, len(items)):
            (x1, y1, x2, y2), ta = items[a]
            (u1, v1, u2, v2), tb = items[c]
            ix, iy = min(x2, u2) - max(x1, u1), min(y2, v2) - max(y1, v1)
            if ix > 0.06 and iy > 0.06 and ("TABLE" in ta or "PICTURE" in ta or
                                            "TABLE" in tb or "PICTURE" in tb):
                ov.append((ta.split()[0], tb.split()[0], round(ix, 2), round(iy, 2)))
    gap = SH - 0.62 - bot
    if gap > 0.95 or over or ov or spill:
        bad.append((i, round(bot, 2), round(gap, 2), over[:2], ov[:2], spill[:3]))

if not bad:
    print("✅ %d 页全部通过（出界 / 留白 / 重叠 / 溢出）" % len(p.slides))
else:
    print("需调整 %d 页：" % len(bad))
    for i, bot, gap, over, ov, spill in bad:
        print("  第 %2d 页 下沿 %.2f 留白 %+.2f" % (i, bot, gap))
        if over:  print("        出界 %s" % (over,))
        if ov:    print("        重叠 %s" % (ov,))
        if spill: print("        溢出 %s" % (spill,))
